"""本地知识库同步主程序。"""

import base64
import io
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Optional
from urllib.parse import quote

import config

# 关掉 pypdf 的 "Ignoring wrong pointing object" 类警告 —— 教科书 PDF 内部 xref
# 表常有非标准条目，pypdf 自动忽略不影响输出，但默认会刷屏污染 Terminal
logging.getLogger("pypdf").setLevel(logging.ERROR)

INVALID_CHARS_RE = re.compile(r'[/\\:*?"<>|^#\[\]%]')
IMG_REF_RE = re.compile(r"(!\[[^\]]*\]\()images/([^)]+)(\))")


def safe_stem(name: str) -> str:
    """规范文件名 stem（不含扩展名），让它在 Obsidian 里安全。"""
    name = INVALID_CHARS_RE.sub("_", name)
    name = re.sub(r"\s+", " ", name)
    name = name.strip(" .")
    return name


def rewrite_image_paths(md_content: str, attachments_subdir: str) -> str:
    """把 ![](images/xxx) 改写为 ![](attachments/<subdir>/xxx)，对 subdir URL-encode。"""
    encoded = quote(attachments_subdir)

    def replace(match: re.Match) -> str:
        return f"{match.group(1)}attachments/{encoded}/{match.group(2)}{match.group(3)}"

    return IMG_REF_RE.sub(replace, md_content)


def compute_target_path(src: Path, source_dir: Path, target_dir: Path) -> Path:
    """02 路径映射到 01 路径，stem 经安全化，后缀改 .md（MinerU 转换路径）。"""
    rel = src.relative_to(source_dir)
    safe_name = safe_stem(rel.stem)
    return target_dir / rel.parent / f"{safe_name}.md"


def is_sidecar(path: Path) -> bool:
    """该后缀的文件不转内容，只在 01 写一个纯 metadata 的 sidecar MD。"""
    return path.suffix.lower() in config.SIDECAR_EXTS


def route(path: Path) -> str:
    """把候选文件分到一条处理路径上。返回：
      'xlsx'    → openpyxl 双读（值+公式，本地）
      'csv'     → csv 模块转 Markdown 表格（本地）
      'sidecar' → 纯 metadata sidecar（老 .xls）
      'pdf'     → pymupdf4llm 本地优先（稀疏 fallback MinerU）
      'pandoc'  → pandoc（docx/rtf/odt/epub，本地）
      'pptx'    → python-pptx + 可选 OCR（本地）
      'passthrough' → 已是文本/Markdown，直接拷贝（仅补 frontmatter）
      'code'    → 代码/结构化文本，包进代码块（本地）
      'mineru'  → MinerU 云（老 .doc/.ppt、.html、图片、PDF 扫描件 fallback）
    """
    ext = path.suffix.lower()
    if ext in config.XLSX_LOCAL_EXTS:
        return "xlsx"
    if ext in config.CSV_EXTS:
        return "csv"
    if ext in config.SIDECAR_EXTS:
        return "sidecar"
    if ext == ".pdf":
        return "pdf"
    if ext in config.PANDOC_EXTS:
        return "pandoc"
    if ext in config.PPTX_LOCAL_EXTS:
        return "pptx"
    if ext in config.PASSTHROUGH_EXTS:
        return "passthrough"
    if ext in config.CODE_EXTS:
        return "code"
    return "mineru"


def find_new_files(
    candidates: list[Path], source_dir: Path, target_dir: Path
) -> list[Path]:
    """转换候选中目标 MD 还不存在的那部分。"""
    return [
        src
        for src in candidates
        if not compute_target_path(src, source_dir, target_dir).exists()
    ]


def is_tool_generated(md_path: Path) -> bool:
    """检查 MD 头部 frontmatter 的 converted_by 是否匹配任一 KNOWN_CONVERTERS。"""
    try:
        with open(md_path, "r", encoding="utf-8") as f:
            head = f.read(2000)
    except OSError:
        return False
    if not head.startswith("---"):
        return False
    return any(
        f"converted_by: {marker}" in head for marker in config.KNOWN_CONVERTERS
    )


def has_source(md_path: Path, target_dir: Path, source_dir: Path) -> bool:
    """
    判断 source_dir 里是否存在一个文件，其 safe_stem 等于 md_path 的 stem，
    且在 source_dir 内的相对父目录与 md_path 在 target_dir 内的相对父目录一致。
    """
    rel = md_path.relative_to(target_dir)
    md_stem = rel.stem
    src_parent = source_dir / rel.parent
    if not src_parent.is_dir():
        return False
    for f in src_parent.iterdir():
        if not f.is_file():
            continue
        if f.suffix.lower() not in config.SUPPORTED_EXTS:
            continue
        if safe_stem(f.stem) == md_stem:
            return True
    return False


def find_orphaned_mds(target_dir: Path, source_dir: Path) -> list[Path]:
    """target_dir 里所有「工具生成 + 无源」的 MD。"""
    orphans: list[Path] = []
    if not target_dir.exists():
        return orphans
    for md_path in sorted(target_dir.rglob("*.md")):
        if not is_tool_generated(md_path):
            continue
        if has_source(md_path, target_dir, source_dir):
            continue
        orphans.append(md_path)
    return orphans


def get_pdf_page_count(path: Path) -> Optional[int]:
    """读 PDF 页数，失败返回 None（不抛）。"""
    try:
        from pypdf import PdfReader
        return len(PdfReader(path).pages)
    except Exception:
        return None


def scan_source(source_dir: Path) -> tuple[list[Path], list[tuple[Path, str]]]:
    """
    遍历 source_dir，返回 (候选文件, 跳过的文件)。
    跳过原因：隐藏文件、不支持的扩展名、超出大小限制。
    """
    candidates: list[Path] = []
    skipped: list[tuple[Path, str]] = []
    if not source_dir.exists():
        return candidates, skipped
    for path in sorted(source_dir.rglob("*")):
        if not path.is_file():
            continue
        if path.name.startswith("."):
            skipped.append((path, "hidden file"))
            continue
        ext = path.suffix.lower()
        if ext not in config.SUPPORTED_EXTS:
            skipped.append((path, f"unsupported: {ext or 'no extension'}"))
            continue
        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > config.MAX_FILE_SIZE_MB:
            skipped.append((path, f"exceeds {config.MAX_FILE_SIZE_MB} MB"))
            continue
        candidates.append(path)
    return candidates, skipped


def build_frontmatter(
    source_filename: str,
    source_type: str,
    converted_at: str,
    pages: Optional[int],
    batch_id: str,
    rel_source_path: Path,
    converted_by: Optional[str] = None,
) -> str:
    """生成 YAML frontmatter 块（含尾部空行）。
    converted_by 默认 'MinerU vlm'（向后兼容 MinerU 流），pymupdf4llm 流传 'pymupdf4llm'。
    """
    if converted_by is None:
        converted_by = f"{config.TOOL_MARKER} vlm"
    lines = [
        "---",
        f'source: "[[{rel_source_path.as_posix()}]]"',
        f'source_filename: "{source_filename}"',
        f"source_type: {source_type}",
        f"converted_at: {converted_at}",
        f"converted_by: {converted_by}",
    ]
    if pages is not None:
        lines.append(f"pages: {pages}")
    lines.append(f"batch_id: {batch_id}")
    lines.append("tags:")
    for tag in config.DEFAULT_TAGS:
        lines.append(f"  - {tag}")
    lines.append("---")
    lines.append("")
    lines.append("")
    return "\n".join(lines)


def _is_cruft(name: str) -> bool:
    """OS 残留文件（.DS_Store / AppleDouble ._* / Thumbs.db），可随空目录一起删。"""
    return name in config.PRUNABLE_CRUFT_NAMES or name.startswith("._")


def _remove_dir_if_cruft_only(d: Path) -> bool:
    """d 为空、或只剩 OS 残留文件时，清掉残留并 rmdir，返回 True；否则不动返回 False。"""
    try:
        children = list(d.iterdir())
    except OSError:
        return False
    if any(c.is_dir() or not _is_cruft(c.name) for c in children):
        return False
    for c in children:
        c.unlink()
    d.rmdir()
    return True


def stage_orphans(
    orphans: list[Path],
    target_dir: Path,
    orphan_root: Path,
    today: str,
) -> list[tuple[Path, Path]]:
    """
    把孤儿 MD 和它的 attachments/<stem>/ 子目录移到 orphan_root/<today>/<rel>。
    返回 [(原位置, 新位置), ...]。
    搬走后向上剪掉空目录（含空的 attachments/），但不删 target_dir 本身。
    """
    moved: list[tuple[Path, Path]] = []
    today_dir = orphan_root / today
    target_resolved = target_dir.resolve()
    for md_path in orphans:
        rel = md_path.relative_to(target_dir)
        dest_md = today_dir / rel
        dest_md.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(md_path), str(dest_md))
        moved.append((md_path, dest_md))

        attach_dir = md_path.parent / "attachments" / rel.stem
        if attach_dir.is_dir():
            dest_attach = today_dir / rel.parent / "attachments" / rel.stem
            dest_attach.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(attach_dir), str(dest_attach))
            moved.append((attach_dir, dest_attach))

        # 剪枝：若 md.parent / "attachments" 空（或只剩 .DS_Store 等残留）则删；
        # 然后从 md.parent 一路向上删空目录，止于 target_dir（不含）。
        attach_parent = md_path.parent / "attachments"
        if attach_parent.is_dir():
            _remove_dir_if_cruft_only(attach_parent)

        cur = md_path.parent
        while cur.resolve() != target_resolved and cur.is_dir():
            if not _remove_dir_if_cruft_only(cur):
                break
            cur = cur.parent
    return moved


def prune_empty_dirs(target_dir: Path) -> list[Path]:
    """Bottom-up rmdir target_dir 下的空子目录；不删 target_dir 自身。

    跳过：路径任一段以 '.' 开头（.obsidian / .claudian / ...），或任一段
    出现在 config.PROTECTED_TARGET_DIR_NAMES（如 '索引'）。返回被删目录列表。
    """
    removed: list[Path] = []
    if not target_dir.is_dir():
        return removed
    all_dirs = sorted(
        (p for p in target_dir.rglob("*") if p.is_dir()),
        key=lambda p: len(p.parts),
        reverse=True,
    )
    for d in all_dirs:
        rel_parts = d.relative_to(target_dir).parts
        if any(part.startswith(".") or part in config.PROTECTED_TARGET_DIR_NAMES
               for part in rel_parts):
            continue
        if _remove_dir_if_cruft_only(d):
            removed.append(d)
    return removed


def extract_zip(
    zip_bytes: bytes,
    target_md_path: Path,
    safe_name: str,
) -> tuple[str, int]:
    """
    解析 MinerU 产物 ZIP：
      - 读出 full.md 内容并返回
      - 把 images/* 拷贝到 target_md_path.parent / 'attachments' / safe_name /
    返回 (md 内容, 拷贝的图片数)。
    """
    with zipfile.ZipFile(io.BytesIO(zip_bytes), "r") as zf:
        names = zf.namelist()
        md_name = next(
            (n for n in names if n.endswith("full.md")),
            None,
        )
        if md_name is None:
            raise ValueError("No full.md in MinerU ZIP")
        md_content = zf.read(md_name).decode("utf-8")

        attach_dir = target_md_path.parent / "attachments" / safe_name
        num_images = 0
        for entry in names:
            if entry == md_name or entry.endswith("/"):
                continue
            if "images/" not in entry:
                continue
            img_name = Path(entry).name
            if not img_name:
                continue
            attach_dir.mkdir(parents=True, exist_ok=True)
            with zf.open(entry) as src, open(attach_dir / img_name, "wb") as dst:
                shutil.copyfileobj(src, dst)
            num_images += 1
        return md_content, num_images


def parse_jwt_exp(token: str) -> Optional[int]:
    """从 JWT 解出 exp（Unix 秒），失败返回 None。不校验签名。"""
    try:
        parts = token.split(".")
        if len(parts) != 3:
            return None
        payload_b64 = parts[1] + "=" * (-len(parts[1]) % 4)
        payload = json.loads(base64.urlsafe_b64decode(payload_b64).decode("utf-8"))
        return payload.get("exp")
    except (ValueError, json.JSONDecodeError, UnicodeDecodeError):
        return None


def validate_token(token: str) -> tuple:
    """返回 (是否有效, 提示信息, 剩余天数)。"""
    exp = parse_jwt_exp(token)
    if exp is None:
        return False, "Token unparseable", 0
    now = int(time.time())
    days = (exp - now) // 86400
    if exp <= now:
        return False, f"Token expired {-days} days ago", days
    return True, f"Token valid, {days} days remaining", days


def setup_logging(log_dir: Path, timestamp: str) -> logging.Logger:
    """同时写文件 logs/sync-<ts>.log 和标准输出。"""
    log_dir.mkdir(parents=True, exist_ok=True)
    log_file = log_dir / f"sync-{timestamp}.log"
    logger = logging.getLogger("sync")
    logger.handlers.clear()
    logger.setLevel(logging.INFO)
    fmt = logging.Formatter("[%(asctime)s] %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
    fh = logging.FileHandler(log_file, encoding="utf-8")
    fh.setFormatter(fmt)
    sh = logging.StreamHandler(sys.stdout)
    sh.setFormatter(fmt)
    logger.addHandler(fh)
    logger.addHandler(sh)
    return logger


def _rel_source_path(src: Path, target_md: Path) -> Path:
    """从 target_md 出发跨过 N 层目录回到项目根，再进 02 拼到原文件。"""
    rel_md = target_md.relative_to(config.TARGET_DIR)
    levels_up = len(rel_md.parts)
    return Path(*([".."] * levels_up)) / config.SOURCE_DIR.name / src.relative_to(config.SOURCE_DIR)


def process_result(src: Path, result, logger: Optional[logging.Logger] = None) -> None:
    """MinerU 单文件成功转换后的落地：解压 → 改写图片路径 → 写 MD（含 frontmatter）。"""
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    target_md.parent.mkdir(parents=True, exist_ok=True)
    safe_name = target_md.stem

    md_content, n_images = extract_zip(result.zip_bytes, target_md, safe_name)
    md_content = rewrite_image_paths(md_content, safe_name)

    fm = build_frontmatter(
        source_filename=src.name,
        source_type=src.suffix.lstrip(".").lower(),
        converted_at=datetime.now().strftime("%Y-%m-%d"),
        pages=result.pages,
        batch_id=result.batch_id,
        rel_source_path=_rel_source_path(src, target_md),
    )
    target_md.write_text(fm + md_content, encoding="utf-8")
    if logger:
        logger.info(f"  ✓ {src.name} → {target_md} ({n_images} images)")
    enrich_frontmatter(target_md, src, logger)


def convert_with_pymupdf4llm(
    src: Path, logger: Optional[logging.Logger] = None
) -> tuple[str, int, Optional[Path]]:
    """用 pymupdf4llm 把 PDF 转成 Markdown 字符串。返回 (md_content, pages, img_dir)。

    本地、秒级、无 quota，但不做 OCR、复杂公式/表格略损。
    开 write_images：把页面里够大的图（≥ PYMUPDF4LLM_IMAGE_SIZE_LIMIT，默认页面 5%）
    抽到一个临时 img_dir，md 里以绝对路径引用；后续 process_pymupdf4llm_result
    把这些图搬进 attachments/ 并改写引用。无图时 img_dir 仍返回（空目录），调用方负责清理。
    """
    import pymupdf4llm
    pages = get_pdf_page_count(src) or 0
    size_mb = src.stat().st_size / (1024 * 1024)
    if logger:
        logger.info(
            f"  → converting {src.name} ({pages} pages, {size_mb:.1f} MB) via pymupdf4llm..."
        )
    t0 = time.time()
    img_dir = Path(tempfile.mkdtemp(prefix="lv_pdfimg_"))
    md_content = pymupdf4llm.to_markdown(
        str(src),
        write_images=True,
        image_path=str(img_dir),
        image_format="png",
        image_size_limit=config.PYMUPDF4LLM_IMAGE_SIZE_LIMIT,
    )
    if logger:
        logger.info(
            f"  ✓ {src.name} (pymupdf4llm, {pages} pages, {time.time() - t0:.1f}s)"
        )
    return md_content, pages, img_dir


def _land_pymupdf4llm_images(
    md_content: str, img_dir: Optional[Path], target_md: Path
) -> tuple[str, int]:
    """把 pymupdf4llm 抽到 img_dir 的图搬进 target_md 旁的 attachments/<stem>/，
    重命名为 ascii 安全名（img-0.png…），并把 md 里的绝对路径引用改写成相对引用。
    返回 (改写后的 md, 落地图片数)。img_dir 用完由调用方清理。

    为什么重命名：pymupdf4llm 默认用「<原文件名>-页-序号.png」，源文件名含空格 / 中文
    会破坏 Markdown 图片链接；统一改成 img-N.png 规避。子目录名 URL-encode，与
    rewrite_image_paths 对 MinerU 的处理保持一致。
    """
    if not img_dir or not Path(img_dir).is_dir():
        return md_content, 0
    files = [p for p in sorted(Path(img_dir).iterdir()) if p.is_file()]
    if not files:
        return md_content, 0
    safe_name = target_md.stem
    attach_dir = target_md.parent / "attachments" / safe_name
    attach_dir.mkdir(parents=True, exist_ok=True)
    enc_sub = quote(safe_name)
    n = 0
    for i, f in enumerate(files):
        clean = f"img-{i}{f.suffix.lower()}"
        shutil.copy2(f, attach_dir / clean)
        # pymupdf4llm 写进 md 的引用 = ](<img_dir>/<原名>)；按原名精确替换
        md_content = md_content.replace(
            f"]({img_dir}/{f.name})", f"](attachments/{enc_sub}/{clean})"
        )
        n += 1
    return md_content, n


def process_pymupdf4llm_result(
    src: Path,
    md_content: str,
    pages: int,
    img_dir: Optional[Path] = None,
    logger: Optional[logging.Logger] = None,
) -> None:
    """pymupdf4llm 转换结果落地：把抽出的图搬进 attachments/，写 MD（含 frontmatter）。"""
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    target_md.parent.mkdir(parents=True, exist_ok=True)
    try:
        md_content, n_images = _land_pymupdf4llm_images(md_content, img_dir, target_md)
    finally:
        if img_dir:
            shutil.rmtree(img_dir, ignore_errors=True)
    fm = build_frontmatter(
        source_filename=src.name,
        source_type=src.suffix.lstrip(".").lower(),
        converted_at=datetime.now().strftime("%Y-%m-%d"),
        pages=pages,
        batch_id="local-pymupdf4llm",
        rel_source_path=_rel_source_path(src, target_md),
        converted_by="pymupdf4llm",
    )
    target_md.write_text(fm + md_content, encoding="utf-8")
    if logger:
        suffix = f"（含 {n_images} 张图）" if n_images else ""
        logger.info(f"  → {target_md}{suffix}")
    enrich_frontmatter(target_md, src, logger)


def extract_excel_metadata(xlsx_path: Path) -> dict:
    """读 xlsx 拿 sheet 名 + 每个 sheet 前 10 行作为样本（给 enrich 提供线索）。
    .xls 老格式 openpyxl 读不了，调用方负责兜底。
    """
    from openpyxl import load_workbook
    wb = load_workbook(xlsx_path, read_only=True, data_only=True)
    sheets = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[tuple] = []
            for row in ws.iter_rows(max_row=10, max_col=12, values_only=True):
                rows.append(row)
            sheets.append({"name": sheet_name, "sample_rows": rows})
    finally:
        wb.close()
    return {"sheets": sheets}


def process_excel_sidecar(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """生成 Excel 的 sidecar MD（纯 metadata，无内容拷贝、无 symlink）。

    01 里只有一个 <stem>.md：含 frontmatter（source 双链回 02 的 xlsx）+ sheet 列表。
    用户想看真表 → 点 frontmatter 里 source 双链 → Obsidian 用 sheet-plus 打开 02 文件。
    Claude 检索时 grep 这个 .md 找得到。
    """
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    target_md.parent.mkdir(parents=True, exist_ok=True)

    sheets_block = ""
    try:
        meta = extract_excel_metadata(src)
        if meta["sheets"]:
            sheets_block = "\n## Sheets\n\n"
            for s in meta["sheets"]:
                sheets_block += f"- **{s['name']}**\n"
    except Exception as exc:
        if logger:
            logger.warning(f"  ⚠ {src.name}: 无法读取 sheet metadata ({exc}); sidecar 仍写出")

    fm = build_frontmatter(
        source_filename=src.name,
        source_type=src.suffix.lstrip(".").lower(),
        converted_at=datetime.now().strftime("%Y-%m-%d"),
        pages=None,
        batch_id="excel-sidecar",
        rel_source_path=_rel_source_path(src, target_md),
        converted_by="excel-sidecar",
    )
    body = (
        f"# {src.stem}\n\n"
        f"Excel 文件的描述索引（sidecar）。**原始表格不在这里**——"
        f"点 frontmatter 顶端的 source 双链跳到 02 里用 sheet-plus 等插件查看。\n"
        f"{sheets_block}"
    )
    target_md.write_text(fm + body, encoding="utf-8")
    if logger:
        logger.info(f"  📋 sidecar {src.name} → {target_md}")
    enrich_frontmatter(target_md, src, logger)


def _write_local_md(
    src: Path,
    target_md: Path,
    body: str,
    *,
    converted_by: str,
    batch_id: str,
    pages: Optional[int],
    logger: Optional[logging.Logger] = None,
    icon: str = "→",
    suffix_msg: str = "",
) -> None:
    """本地转换器共用的落地：写 frontmatter + body，然后 enrich。"""
    target_md.parent.mkdir(parents=True, exist_ok=True)
    fm = build_frontmatter(
        source_filename=src.name,
        source_type=src.suffix.lstrip(".").lower(),
        converted_at=datetime.now().strftime("%Y-%m-%d"),
        pages=pages,
        batch_id=batch_id,
        rel_source_path=_rel_source_path(src, target_md),
        converted_by=converted_by,
    )
    target_md.write_text(fm + body, encoding="utf-8")
    if logger:
        logger.info(f"  {icon} {src.name} → {target_md}{suffix_msg}")
    enrich_frontmatter(target_md, src, logger)


# ── Excel：openpyxl 双读（值 + 公式）────────────────────────────────
def convert_xlsx(src: Path) -> str:
    """openpyxl 双读 .xlsx → Markdown body。
    - data_only=True  拿 Excel **缓存的**计算值（真实 Excel 存的文件才有缓存）
    - data_only=False 拿公式串（'=SUM(...)'）
    每个 sheet 输出带坐标的值表（列字母 + 行号）+ 一份公式清单，方便 LLM 顺坐标读懂模型。
    单 sheet 超 EXCEL_MAX_CELLS_PER_SHEET 截断；列跨度过宽时退化为 key-value 列表。
    """
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    wb_v = load_workbook(src, read_only=True, data_only=True)
    wb_f = load_workbook(src, read_only=True, data_only=False)
    parts = [f"# {src.stem}\n"]
    try:
        for name in wb_v.sheetnames:
            ws_v = wb_v[name]
            ws_f = wb_f[name]
            parts.append(f"\n## {name}\n")

            cells: dict[tuple[int, int], str] = {}
            formulas: list[tuple[str, str]] = []
            min_r = min_c = None
            max_r = max_c = 0
            seen = 0
            truncated = False
            for v_row, f_row in zip(ws_v.iter_rows(), ws_f.iter_rows()):
                for v_cell, f_cell in zip(v_row, f_row):
                    val = v_cell.value
                    fval = f_cell.value
                    if val is None and fval is None:
                        continue
                    if seen >= config.EXCEL_MAX_CELLS_PER_SHEET:
                        truncated = True
                        break
                    seen += 1
                    r, c = v_cell.row, v_cell.column
                    coord = f"{get_column_letter(c)}{r}"
                    if isinstance(fval, str) and fval.startswith("="):
                        formulas.append((coord, fval))
                    disp = "" if val is None else str(val).replace("\n", " ").replace("|", "\\|")
                    cells[(r, c)] = disp
                    min_r = r if min_r is None else min(min_r, r)
                    max_r = max(max_r, r)
                    min_c = c if min_c is None else min(min_c, c)
                    max_c = max(max_c, c)
                if truncated:
                    break

            if not cells:
                parts.append("_（空 sheet）_\n")
                continue

            width = max_c - min_c + 1
            if width <= 64:
                # 网格表：表头是列字母，第一列是行号
                header = ["", *[get_column_letter(c) for c in range(min_c, max_c + 1)]]
                lines = [
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join("---" for _ in header) + " |",
                ]
                for r in range(min_r, max_r + 1):
                    row = [str(r)] + [cells.get((r, c), "") for c in range(min_c, max_c + 1)]
                    lines.append("| " + " | ".join(row) + " |")
                parts.append("\n".join(lines) + "\n")
            else:
                # 跨度过宽 → key-value 列表（避免巨宽空网格）
                kv = [f"- `{get_column_letter(c)}{r}`: {cells[(r, c)]}"
                      for (r, c) in sorted(cells)]
                parts.append("\n".join(kv) + "\n")

            if formulas:
                parts.append("\n**公式**：\n")
                parts.append("\n".join(f"- `{coord}`: `{f}`" for coord, f in formulas) + "\n")
            if truncated:
                parts.append(
                    f"\n> ⚠️ 此 sheet 超过 {config.EXCEL_MAX_CELLS_PER_SHEET} 个非空单元格，已截断。"
                    f"完整数据请点 frontmatter 的 source 双链回 02 原表查看。\n"
                )
    finally:
        wb_v.close()
        wb_f.close()
    return "\n".join(parts) + "\n"


def process_xlsx(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """.xlsx 本地转换落地（值 + 公式）。"""
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    body = convert_xlsx(src)
    _write_local_md(
        src, target_md, body,
        converted_by="excel-openpyxl",
        batch_id="local-openpyxl",
        pages=None,
        logger=logger,
        icon="📊",
    )


# ── CSV / TSV：csv 模块 → Markdown 表格 ───────────────────────────
def convert_csv(src: Path) -> str:
    """.csv/.tsv → Markdown 表格 body。超 CSV_MAX_ROWS 行截断 + 提示。"""
    import csv
    delim = "\t" if src.suffix.lower() == ".tsv" else ","
    rows: list[list[str]] = []
    with open(src, newline="", encoding="utf-8-sig") as f:
        for row in csv.reader(f, delimiter=delim):
            rows.append(row)

    parts = [f"# {src.stem}\n"]
    if not rows:
        return "\n".join(parts) + "\n_（空文件）_\n"

    truncated = len(rows) > config.CSV_MAX_ROWS
    rows = rows[: config.CSV_MAX_ROWS]
    ncols = max(len(r) for r in rows)

    def esc(c: str) -> str:
        return str(c).replace("\n", " ").replace("|", "\\|")

    def pad(r: list[str]) -> list[str]:
        return [esc(c) for c in r] + [""] * (ncols - len(r))

    lines = [
        "| " + " | ".join(pad(rows[0])) + " |",
        "| " + " | ".join("---" for _ in range(ncols)) + " |",
    ]
    for r in rows[1:]:
        lines.append("| " + " | ".join(pad(r)) + " |")
    parts.append("\n".join(lines) + "\n")
    if truncated:
        parts.append(
            f"\n> ⚠️ 仅显示前 {config.CSV_MAX_ROWS} 行；完整数据见 frontmatter 的 source 双链回 02 原文件。\n"
        )
    return "\n".join(parts) + "\n"


def process_csv(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """.csv/.tsv 本地转 Markdown 表格落地。"""
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    body = convert_csv(src)
    _write_local_md(
        src, target_md, body,
        converted_by="csv",
        batch_id="local-csv",
        pages=None,
        logger=logger,
        icon="📈",
    )


# ── Word/RTF/ODT/EPUB：pandoc → Markdown ──────────────────────────
def convert_via_pandoc(src: Path, target_md: Path) -> tuple[str, int]:
    """pandoc 把 docx/rtf/odt/epub 转成 GitHub-flavored Markdown（按扩展名自动识别输入格式），
    嵌入图片抽到 attachments/<stem>/。返回 (md_content, 图片数)。
    pandoc 失败抛 RuntimeError，调用方决定兜底。
    """
    safe_name = target_md.stem
    attach_dir = target_md.parent / "attachments" / safe_name
    attach_dir.parent.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        [config.PANDOC_CMD, str(src), "-t", "gfm", "--wrap=none",
         f"--extract-media={attach_dir}"],
        capture_output=True, text=True, timeout=config.PANDOC_TIMEOUT_SEC,
    )
    if result.returncode != 0:
        raise RuntimeError(f"pandoc exit {result.returncode}: {result.stderr.strip()[:200]}")
    md = result.stdout
    # pandoc 把图片引用写成 <attach_dir>/media/x.png（绝对路径）。改写成相对路径。
    # 同时覆盖 pandoc 可能对路径做的 %20 编码。
    rel_prefix = f"attachments/{quote(safe_name)}/"
    md = md.replace(str(attach_dir) + "/", rel_prefix)
    md = md.replace(quote(str(attach_dir)) + "/", rel_prefix)
    n_images = 0
    media = attach_dir / "media"
    if media.is_dir():
        n_images = sum(1 for p in media.iterdir() if p.is_file())
    return md, n_images


def process_pandoc(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """docx/rtf/odt/epub 本地转换落地（pandoc）。"""
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    md, n_images = convert_via_pandoc(src, target_md)
    _write_local_md(
        src, target_md, md,
        converted_by="pandoc",
        batch_id="local-pandoc",
        pages=None,
        logger=logger,
        icon="📝",
        suffix_msg=f" ({n_images} images)" if n_images else "",
    )


# ── PowerPoint：python-pptx（文字/表格/备注）+ 可选图片 OCR ────────
def _iter_pptx_shapes(shapes):
    """递归遍历 shapes（展开 group），逐个 yield 叶子 shape。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def _pptx_table_to_md(table) -> str:
    """python-pptx table → Markdown 表格。"""
    rows = []
    for row in table.rows:
        rows.append([c.text.strip().replace("\n", " ").replace("|", "\\|") for c in row.cells])
    if not rows:
        return ""
    header = rows[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _pptx_chart_to_md(chart) -> str:
    """python-pptx **原生图表对象** → Markdown 表（类别 + 各 series 数值）。
    图表既不是 table 也不是 picture，不专门提取就会整个丢掉里面的数字。
    """
    def esc(v) -> str:
        return str("" if v is None else v).replace("\n", " ").replace("|", "\\|")

    series = list(chart.series)
    if not series:
        return ""
    try:
        cats = [esc(c) for c in chart.plots[0].categories]
    except (IndexError, ValueError, AttributeError):
        cats = []
    cols = [list(s.values) for s in series]
    nrows = max([len(cats)] + [len(c) for c in cols])

    header = ["类别", *[esc(s.name) for s in series]]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for i in range(nrows):
        cat = cats[i] if i < len(cats) else str(i + 1)
        row = [cat] + [esc(col[i]) if i < len(col) else "" for col in cols]
        lines.append("| " + " | ".join(row) + " |")

    title = ""
    try:
        if chart.has_title and chart.chart_title.text_frame.text.strip():
            title = f"**图表：{esc(chart.chart_title.text_frame.text.strip())}**\n\n"
    except Exception:
        pass
    return title + "\n".join(lines)


def ocr_image_via_claude(img_path: Path, logger: Optional[logging.Logger] = None) -> str:
    """用 claude -p 的 Read 工具 OCR 一张图片，返回转写文本。失败/空返回 ""（不抛）。
    复用用户 Claude Code 订阅，无需 ANTHROPIC_API_KEY、不烧 MinerU quota。
    """
    if shutil.which(config.ENRICH_CLAUDE_CMD) is None:
        return ""
    abspath = str(img_path.resolve())
    prompt = (
        f"Use your Read tool to open the image at {abspath} and transcribe ALL text, "
        f"numbers and data you see, preserving structure (render any tables as Markdown). "
        f"Output ONLY the transcription, no commentary. If there is no readable text, "
        f"output nothing."
    )
    try:
        result = subprocess.run(
            [config.ENRICH_CLAUDE_CMD, "-p", prompt],
            capture_output=True, text=True, timeout=config.OCR_TIMEOUT_SEC,
        )
    except subprocess.TimeoutExpired:
        if logger:
            logger.warning(f"  ⚠ OCR {img_path.name} 超时（{config.OCR_TIMEOUT_SEC}s）")
        return ""
    except Exception as exc:
        if logger:
            logger.warning(f"  ⚠ OCR {img_path.name} 调用异常: {exc}")
        return ""
    if result.returncode != 0:
        if logger:
            logger.warning(f"  ⚠ OCR {img_path.name} 失败 (exit {result.returncode})")
        return ""
    return result.stdout.strip()


def convert_pptx(
    src: Path, target_md: Path, logger: Optional[logging.Logger] = None
) -> tuple[str, int, int]:
    """python-pptx 抽取 .pptx → Markdown body。
    每页：标题/正文文字 + 表格 + 嵌入图片（抽到 attachments/<stem>/，可选 OCR）+ 备注。

    「智能 OCR」（见 config）：图片按内容去重只 OCR 一次、跳过装饰小图、唯一内容图并发跑，
    避免每张装饰 logo 都串行起一个 claude 子进程（慢的根因）。
    返回 (md_content, slide 数, 图片引用数)。
    """
    import hashlib
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(src))
    safe_name = target_md.stem
    attach_dir = target_md.parent / "attachments" / safe_name

    # ── Pass 1：抽文字/表格/图片，建中间表示；图片按 blob 去重 ──────
    slides_repr: list[list[tuple]] = []   # 每页一串 ("md", text) / ("img", hash, name, rel)
    by_hash: dict[str, str] = {}          # blob_hash → img_name（去重抽取，只写一次）
    ocr_targets: dict[str, Path] = {}     # 需要 OCR 的唯一图：blob_hash → 文件路径
    n_images = 0
    n_skipped_small = 0

    for slide in prs.slides:
        elems: list[tuple] = []
        for shape in _iter_pptx_shapes(slide.shapes):
            if getattr(shape, "has_chart", False):
                cm = _pptx_chart_to_md(shape.chart)
                if cm:
                    elems.append(("md", cm))
                continue
            if shape.has_table:
                tbl = _pptx_table_to_md(shape.table)
                if tbl:
                    elems.append(("md", tbl))
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                elems.append(("md", shape.text_frame.text.strip()))
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    blob = img.blob
                except Exception:
                    continue
                n_images += 1
                h = hashlib.sha1(blob).hexdigest()
                if h not in by_hash:
                    ext = (img.ext or "png").lstrip(".")
                    img_name = f"img-{h[:12]}.{ext}"
                    attach_dir.mkdir(parents=True, exist_ok=True)
                    (attach_dir / img_name).write_bytes(blob)
                    by_hash[h] = img_name
                    if config.OCR_PPTX_IMAGES:
                        if len(blob) >= config.OCR_MIN_IMAGE_BYTES:
                            ocr_targets[h] = attach_dir / img_name
                        else:
                            n_skipped_small += 1
                img_name = by_hash[h]
                rel = f"attachments/{quote(safe_name)}/{img_name}"
                elems.append(("img", h, img_name, rel))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                elems.append(("md", f"**备注（Notes）**：{notes}"))
        slides_repr.append(elems)

    # ── 唯一内容图并发 OCR ──────────────────────────────────────────
    ocr_map: dict[str, str] = {}
    if ocr_targets:
        from concurrent.futures import ThreadPoolExecutor
        if logger:
            logger.info(
                f"    OCR {len(ocr_targets)} 张唯一内容图（并发 {config.OCR_MAX_WORKERS}）"
                f"，去重省掉 {len(by_hash) - len(ocr_targets) - n_skipped_small} 张重复 / "
                f"跳过 {n_skipped_small} 张装饰小图..."
            )
        with ThreadPoolExecutor(max_workers=config.OCR_MAX_WORKERS) as ex:
            futs = {ex.submit(ocr_image_via_claude, p, logger): h
                    for h, p in ocr_targets.items()}
            for fut in futs:
                try:
                    ocr_map[futs[fut]] = fut.result()
                except Exception:
                    ocr_map[futs[fut]] = ""

    # ── Pass 2：渲染 ────────────────────────────────────────────────
    parts = [f"# {src.stem}\n"]
    for idx, elems in enumerate(slides_repr, 1):
        parts.append(f"\n## Slide {idx}\n")
        for e in elems:
            if e[0] == "md":
                parts.append(e[1])
            else:
                _, h, img_name, rel = e
                parts.append(f"![{img_name}]({rel})")
                ocr = ocr_map.get(h, "")
                if ocr:
                    quoted = "\n".join(f"> {ln}" for ln in ocr.splitlines())
                    parts.append(f"> **图片文字（OCR）**：\n>\n{quoted}")
    return "\n".join(parts) + "\n", len(slides_repr), n_images


def process_pptx(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """.pptx 本地转换落地（python-pptx + 可选图片 OCR）。"""
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    body, n_slides, n_images = convert_pptx(src, target_md, logger)
    _write_local_md(
        src, target_md, body,
        converted_by="python-pptx",
        batch_id="local-pptx",
        pages=n_slides,
        logger=logger,
        icon="📑",
        suffix_msg=f" ({n_slides} slides, {n_images} images)",
    )


# ── Markdown：直接拷贝（仅补 frontmatter，绝不动正文）──────────────
def _inject_passthrough_keys(raw: str, rel_source_path: Path) -> Optional[str]:
    """往已有 frontmatter 里补 source / converted_by: passthrough（缺哪个补哪个），
    其余字段与正文一字不动。无法定位 frontmatter 闭合 → 返回 None（调用方改前置）。
    """
    if not raw.startswith("---"):
        return None
    closing_pos = raw.find("\n---", 3)  # frontmatter 闭合的 '---'
    if closing_pos == -1:
        return None
    head = raw[:closing_pos]
    inject: list[str] = []
    if "\nsource:" not in head and not head.startswith("source:"):
        inject.append(f'source: "[[{rel_source_path.as_posix()}]]"')
    if "converted_by:" not in head:
        inject.append("converted_by: passthrough")
    if not inject:
        return raw
    return head + "\n" + "\n".join(inject) + raw[closing_pos:]


def process_md_passthrough(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """02 里的 .md/.markdown/.txt 直接拷到 01：
    - 已有 frontmatter → 只补 source + converted_by: passthrough，原 frontmatter/正文不动
    - 无 frontmatter   → 前置最小 frontmatter，正文原样跟随
    然后照常 enrich（只动 frontmatter）。
    """
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    target_md.parent.mkdir(parents=True, exist_ok=True)
    raw = src.read_text(encoding="utf-8", errors="replace")
    rel = _rel_source_path(src, target_md)

    out = _inject_passthrough_keys(raw, rel)
    if out is None:
        fm = build_frontmatter(
            source_filename=src.name,
            source_type=src.suffix.lstrip(".").lower(),
            converted_at=datetime.now().strftime("%Y-%m-%d"),
            pages=None,
            batch_id="local-passthrough",
            rel_source_path=rel,
            converted_by="passthrough",
        )
        out = fm + raw

    target_md.write_text(out, encoding="utf-8")
    if logger:
        logger.info(f"  📄 passthrough {src.name} → {target_md}")
    enrich_frontmatter(target_md, src, logger)


# ── 代码/结构化文本：包进代码块（带语言标注）──────────────────────
def process_code_passthrough(src: Path, logger: Optional[logging.Logger] = None) -> None:
    """02 里的 .json/.yaml/.py/... 包进 ``` 代码块写到 01（带语言标注），补 frontmatter。
    正文（代码原文）不动，只是裹一层 fence 让 Obsidian 渲染 + 标语言。
    """
    target_md = compute_target_path(src, config.SOURCE_DIR, config.TARGET_DIR)
    ext = src.suffix.lower()
    raw = src.read_text(encoding="utf-8", errors="replace")
    lang = config.CODE_FENCE_LANG.get(ext, ext.lstrip("."))
    # 防御：原文若含 ``` 用更长的围栏避免提前闭合
    fence = "````" if "```" in raw else "```"
    body = f"# {src.stem}\n\n{fence}{lang}\n{raw}\n{fence}\n"
    _write_local_md(
        src, target_md, body,
        converted_by="passthrough",
        batch_id="local-code",
        pages=None,
        logger=logger,
        icon="🧩",
    )


def enrich_frontmatter(
    md_path: Path,
    src: Path,
    logger: Optional[logging.Logger] = None,
) -> bool:
    """调 `claude -p` 给已写出的 MD frontmatter 追加 abstract / auto_tags /
    synonyms / key_data。**只动 frontmatter，不动正文**。失败时不抛、不阻塞。
    返回是否成功 enrich。
    """
    if not config.ENRICH_FRONTMATTER:
        return False
    if shutil.which(config.ENRICH_CLAUDE_CMD) is None:
        if logger:
            logger.warning(f"  ⚠ '{config.ENRICH_CLAUDE_CMD}' 不在 PATH，跳过 enrich")
        return False

    try:
        full_text = md_path.read_text(encoding="utf-8")
    except OSError as exc:
        if logger:
            logger.warning(f"  ⚠ 读 {md_path.name} 失败: {exc}")
        return False

    head = full_text[: config.ENRICH_SAMPLE_HEAD_CHARS]
    tail = ""
    if len(full_text) > config.ENRICH_SAMPLE_HEAD_CHARS + config.ENRICH_SAMPLE_TAIL_CHARS:
        tail = full_text[-config.ENRICH_SAMPLE_TAIL_CHARS:]

    prompt = (
        f"Read this Markdown document (extracted from {src.name}) and emit ONLY a YAML "
        f"snippet (no code fences, no commentary, no surrounding ---) with these keys:\n\n"
        f"abstract: |\n"
        f"  Three-sentence summary in the document's primary language.\n"
        f"auto_tags: [tag1, tag2, ...]   # 5-10 specific lowercase tags\n"
        f"synonyms: [term1, ...]          # multi-language synonyms (English + Chinese)\n"
        f"key_data:                        # optional, 3-5 important facts/numbers\n"
        f"  - \"fact 1\"\n\n"
        f"---DOCUMENT HEAD---\n{head}\n"
        + (f"\n---DOCUMENT TAIL---\n{tail}\n" if tail else "")
        + "\n---END---\n\n"
        f"Output ONLY the YAML keys (abstract/auto_tags/synonyms/key_data), nothing else."
    )

    try:
        result = subprocess.run(
            [config.ENRICH_CLAUDE_CMD, "-p", prompt],
            capture_output=True,
            text=True,
            timeout=config.ENRICH_TIMEOUT_SEC,
        )
    except subprocess.TimeoutExpired:
        if logger:
            logger.warning(f"  ⚠ enrich {md_path.name} 超时（{config.ENRICH_TIMEOUT_SEC}s）")
        return False
    except Exception as exc:
        if logger:
            logger.warning(f"  ⚠ enrich {md_path.name} 调用异常: {exc}")
        return False

    if result.returncode != 0:
        if logger:
            logger.warning(
                f"  ⚠ enrich {md_path.name} 失败 (exit {result.returncode}): "
                f"{result.stderr.strip()[:200]}"
            )
        return False

    yaml_snippet = result.stdout.strip()
    # 防御：去掉模型可能误加的 ``` 围栏
    if yaml_snippet.startswith("```"):
        yaml_snippet = re.sub(r"^```[\w]*\n?", "", yaml_snippet)
        yaml_snippet = re.sub(r"\n?```$", "", yaml_snippet)
    if not yaml_snippet:
        return False

    # 把 yaml_snippet 插入到 frontmatter 的闭合 `---` 之前
    # frontmatter 形如：---\n<keys>\n---\n\n<body>
    closing_pos = full_text.find("\n---\n", 4)  # 跳过开头 ---
    if closing_pos == -1:
        if logger:
            logger.warning(f"  ⚠ {md_path.name} frontmatter 闭合标记找不到，跳过 enrich")
        return False

    new_text = (
        full_text[: closing_pos + 1]
        + yaml_snippet
        + "\n"
        + full_text[closing_pos + 1 :]
    )
    md_path.write_text(new_text, encoding="utf-8")
    if logger:
        logger.info(f"  ✨ enriched {md_path.name}")
    return True


def report_skipped(
    skipped: list[tuple[Path, str]],
    source_dir: Path,
    logger: logging.Logger,
) -> None:
    """在 sync 末尾醒目列出**未被转换**的文件 + 处理建议。
    隐藏文件（.DS_Store 等）不打扰用户；只报「不支持的类型」和「超大小上限」两类。
    """
    unsupported = [(p, r) for p, r in skipped if r.startswith("unsupported")]
    oversize = [(p, r) for p, r in skipped if "exceeds" in r]
    if not unsupported and not oversize:
        return

    logger.info("")
    logger.info("════════════ ⚠️  有文件未被转换，请注意 ════════════")
    if unsupported:
        from collections import defaultdict
        by_ext: dict[str, list[Path]] = defaultdict(list)
        for p, _ in unsupported:
            by_ext[p.suffix.lower() or "(无扩展名)"].append(p)
        logger.info(f"【不支持的类型】共 {len(unsupported)} 个文件：")
        for ext, paths in sorted(by_ext.items()):
            hint = config.UNSUPPORTED_HINTS.get(ext, config.DEFAULT_UNSUPPORTED_HINT)
            logger.info(f"  • {ext}（{len(paths)} 个）→ {hint}")
            for p in paths:
                logger.info(f"        {p.relative_to(source_dir)}")
    if oversize:
        logger.info(f"【超过大小上限 {config.MAX_FILE_SIZE_MB} MB】共 {len(oversize)} 个文件：")
        for p, _ in oversize:
            logger.info(f"  • {p.relative_to(source_dir)} → 调大 config.MAX_FILE_SIZE_MB 或拆分后再放进 02")
    logger.info("（以上文件留在 02，未生成 01 的 MD；按建议处理后下次 sync 会重新尝试）")
    logger.info("═══════════════════════════════════════════════")
    logger.info("")


def _launcher_body(sync_py: Path) -> str:
    """生成放进用户文件夹的双击 launcher 内容。

    工具本体（sync.py）和用户数据是分离的——插件安装时 sync.py 埋在
    ~/.claude/plugins/cache/.../scripts/ 下，跟原始文件目录完全两处。所以这里
    把 sync.py 的【绝对路径】硬编码进 launcher，不能用相对路径的 cd "$DIR"。
    """
    return (
        "#!/bin/bash\n"
        "# 本地知识库 · 双击同步入口（local-vault 自动生成，勿手改——升级会被覆盖刷新）\n"
        "# 用法：把要转换的文件拖进本文件夹，然后双击我 → 自动转成 Markdown 进 vault。\n"
        f'python3 "{sync_py}"\n'
        "EXIT_CODE=$?\n"
        "echo \"\"\n"
        "echo \"─────── 同步结束（退出码 $EXIT_CODE），按回车关闭窗口 ───────\"\n"
        "read -r\n"
        "exit $EXIT_CODE\n"
    )


# 用来识别「我们自己生成的」launcher（_launcher_body 注释里一定带这句）。
_LAUNCHER_MARK = "local-vault 自动生成"


def _is_our_launcher(text: str) -> bool:
    return _LAUNCHER_MARK in text


def _cleanup_stale_source_launcher(
    source_dir: Path, home: Path, logger: Optional[logging.Logger] = None
) -> None:
    """旧版本把 launcher 放进 SOURCE(02)；现在落点改到父目录（本地知识库根）。
    若 SOURCE 里还残留我们自己生成的 sync.command（且不是新落点本身），删掉以消除重复。
    用户手写的（没有我们的标记）一律不碰。"""
    try:
        stale = source_dir / "sync.command"
        if stale.resolve() == (home / "sync.command").resolve():
            return
        if stale.exists() and _is_our_launcher(stale.read_text(encoding="utf-8")):
            stale.unlink()
            if logger:
                logger.info(f"已清理旧位置的重复双击入口：{stale}")
    except OSError:
        pass


def _ask_launcher_update(launcher: Path) -> bool:
    """终端交互：已有 sync.command 且与最新版不同时，问更新还是跳过。
    返回 True=更新。空回车 / EOF / Ctrl-C 一律按跳过处理。"""
    try:
        ans = input(
            f"\n检测到已存在 sync.command（与最新版不同）：\n  {launcher}\n"
            f"更新它吗？[u]更新 / [s]跳过（默认跳过）: "
        ).strip().lower()
    except (EOFError, KeyboardInterrupt):
        return False
    return ans in {"u", "update", "y", "yes"}


def ensure_clickable_launcher(
    source_dir: Path,
    logger: Optional[logging.Logger] = None,
    *,
    interactive: Optional[bool] = None,
) -> Optional[Path]:
    """在【本地知识库根】（source_dir 的父目录）放/刷新可双击的 sync.command。

    落点是父目录而非 SOURCE 本身：双击入口和「往里拖文件」的原始目录分开；对
    /plugin install（sync.py 埋在 cache、PROJECT_ROOT 不在数据旁）也成立。
    仅 macOS——.command 是 Finder 双击格式；其它平台静默跳过。
    .env 里 KB_NO_LAUNCHER=1 可整体关掉（见 config.INSTALL_CLICKABLE_LAUNCHER）。

    存在性处理：
      - 不存在 → 直接创建。
      - 已存在且与最新版一致 → 不动（返回 None）。
      - 已存在但不同 →
          · 交互终端(TTY)：提示用户「更新 / 跳过」，由用户拍板；
          · 非交互 + 是我们自己生成的旧 launcher：静默自愈（刷新路径，沿用插件升级自愈）；
          · 非交互 + 用户自定义 launcher：保留不动，绝不覆盖。
    返回写入/更新的路径，未写则 None。
    """
    if not config.INSTALL_CLICKABLE_LAUNCHER:
        return None
    if sys.platform != "darwin":
        return None
    try:
        sync_py = (config.TOOL_DIR / "sync.py").resolve()
        desired = _launcher_body(sync_py)
        home = source_dir.parent if source_dir.parent != source_dir else source_dir

        _cleanup_stale_source_launcher(source_dir, home, logger)

        launcher = home / "sync.command"
        if not launcher.exists():
            launcher.write_text(desired, encoding="utf-8")
            launcher.chmod(0o755)
            if logger:
                logger.info(f"已在本地知识库根放置双击入口：{launcher}（双击即可同步）")
            return launcher

        current = launcher.read_text(encoding="utf-8")
        if current == desired:
            return None  # 已是最新，无需动作

        if interactive is None:
            interactive = sys.stdin.isatty() and sys.stdout.isatty()

        if interactive:
            if not _ask_launcher_update(launcher):
                if logger:
                    logger.info(f"保留现有 sync.command（用户选择跳过）：{launcher}")
                return None
        elif not _is_our_launcher(current):
            if logger:
                logger.info(f"检测到自定义 sync.command，保留不动：{launcher}")
            return None
        # 走到这：交互里选了更新，或非交互且是我们自己的旧 launcher（静默自愈）
        launcher.write_text(desired, encoding="utf-8")
        launcher.chmod(0o755)
        if logger:
            logger.info(f"已更新双击入口：{launcher}")
        return launcher
    except OSError as e:
        if logger:
            logger.warning(f"放置 sync.command 失败（不影响命令行使用）：{e}")
        return None


def _write_env_file(env_path: Path, source_dir, target_dir, token: str = "") -> None:
    """把首次运行向导收集到的配置写进 .env（config.load_dotenv 下次会读它）。"""
    lines = [
        "# local-vault 配置（首次运行向导生成）",
        f"KB_SOURCE_DIR={source_dir}",
        f"KB_TARGET_DIR={target_dir}",
    ]
    if token:
        lines.append(f"MINERU_TOKEN={token}")
    env_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def first_run_setup() -> bool:
    """终端交互向导：没配 KB_SOURCE_DIR/KB_TARGET_DIR 时引导用户配置 + 教用法。
    只在 TTY 下调用（见 main）。写好 .env 后需重新运行（config 已按旧值加载）。
    返回是否成功写入配置。
    """
    print("\n" + "═" * 56)
    print("  📚 本地知识库 · 首次配置向导")
    print("═" * 56)
    print("这个工具把原始文件（PDF / Word / Excel / PPT / 图片 / md…）转成")
    print("Markdown，写进一个 vault，方便 Obsidian 阅读 + AI 检索。")
    print("先告诉我两个目录（可用 ~ 开头的家目录路径）：\n")

    try:
        source = input("1) 原始文件目录（你往里拖文件）: ").strip().strip("'\"")
        target = input("2) Markdown vault 目录（产出 .md 放这）: ").strip().strip("'\"")
    except (EOFError, KeyboardInterrupt):
        print("\n已取消。下次运行再配即可。")
        return False
    if not source or not target:
        print("\n⚠️  路径不能为空，已取消。下次运行再配即可。")
        return False

    src = Path(source).expanduser()
    tgt = Path(target).expanduser()
    src.mkdir(parents=True, exist_ok=True)
    tgt.mkdir(parents=True, exist_ok=True)

    print("\n3) MinerU token（可选；只为处理老 .doc/.ppt、扫描件、图片。回车跳过）")
    try:
        token = input("   token: ").strip()
    except (EOFError, KeyboardInterrupt):
        token = ""

    _write_env_file(config.TOOL_DIR / ".env", src, tgt, token)
    launcher = ensure_clickable_launcher(src)

    print("\n" + "─" * 56)
    print("✅ 配置已写入 .env。怎么用：")
    print(f"  1. 把要转换的文件拖进：{src}")
    if launcher:
        print(f"  2. 双击 {launcher} → 自动转 Markdown")
        print("     （它在原始目录的上一层；命令行等价：python3 sync.py）")
    else:
        print("  2. 再运行一次 python3 sync.py → 自动转 Markdown")
    print(f"  3. 产出的 .md 在：{tgt}")
    print("     用 Obsidian 打开它阅读，或在 Claude 里基于它问答")
    print("  · 想重转某个文件：先删掉 vault 里对应的 .md 再运行")
    print("  · 不支持的类型（音视频 / iWork / 压缩包等）运行后会在末尾提示怎么处理")
    print("─" * 56 + "\n")
    return True


def main() -> int:
    from mineru_client import MinerUClient

    timestamp = datetime.now().strftime("%Y-%m-%d-%H%M%S")
    logger = setup_logging(config.LOG_DIR, timestamp)
    logger.info("=== Sync started ===")

    # 注意：MINERU_TOKEN 在真正需要调 MinerU 时才校验（见 Stage 3）。
    # 这样只有 Excel 或纯数字 PDF 的同步即使 token 未设也能跑完。

    # 没配好源/目标目录：终端里走交互向导；非交互（Claude 子进程等）保持报错
    if not config.SOURCE_DIR.exists() or not config.TARGET_DIR.exists():
        if sys.stdin.isatty():
            ok = first_run_setup()
            return 0 if ok else 1
        if not config.SOURCE_DIR.exists():
            logger.error(f"源目录不存在: {config.SOURCE_DIR}")
            logger.error("  → 设环境变量 KB_SOURCE_DIR（或在终端直接运行一次走配置向导）")
        if not config.TARGET_DIR.exists():
            logger.error(f"目标目录不存在: {config.TARGET_DIR}")
            logger.error("  → 设环境变量 KB_TARGET_DIR（或在终端直接运行一次走配置向导）")
        return 1

    # 兜底放置双击入口：Claude 帮配 .env 的路径不走向导，这里幂等补上（macOS）
    ensure_clickable_launcher(config.SOURCE_DIR, logger)

    today = datetime.now().strftime("%Y-%m-%d")

    logger.info(f"扫描 {config.SOURCE_DIR}...")
    candidates, skipped = scan_source(config.SOURCE_DIR)
    logger.info(f"  扫描到 {len(candidates)} 个候选文件，{len(skipped)} 个跳过")
    for path, reason in skipped:
        logger.info(f"  Skipped: {path.relative_to(config.SOURCE_DIR)} ({reason})")

    # ── 路由分流（见 route()）─────────────────────────────────────
    buckets: dict[str, list[Path]] = {
        "xlsx": [], "csv": [], "sidecar": [], "pdf": [], "pandoc": [], "pptx": [],
        "passthrough": [], "code": [], "mineru": [],
    }
    for c in candidates:
        buckets[route(c)].append(c)

    new_xlsx = find_new_files(buckets["xlsx"], config.SOURCE_DIR, config.TARGET_DIR)
    new_csv = find_new_files(buckets["csv"], config.SOURCE_DIR, config.TARGET_DIR)
    new_sidecars = find_new_files(buckets["sidecar"], config.SOURCE_DIR, config.TARGET_DIR)
    new_pdfs = find_new_files(buckets["pdf"], config.SOURCE_DIR, config.TARGET_DIR)
    new_pandoc = find_new_files(buckets["pandoc"], config.SOURCE_DIR, config.TARGET_DIR)
    new_pptx = find_new_files(buckets["pptx"], config.SOURCE_DIR, config.TARGET_DIR)
    new_md = find_new_files(buckets["passthrough"], config.SOURCE_DIR, config.TARGET_DIR)
    new_code = find_new_files(buckets["code"], config.SOURCE_DIR, config.TARGET_DIR)
    new_others = find_new_files(buckets["mineru"], config.SOURCE_DIR, config.TARGET_DIR)

    orphans = find_orphaned_mds(config.TARGET_DIR, config.SOURCE_DIR)
    logger.info(
        f"  待转 — xlsx: {len(new_xlsx)}, csv: {len(new_csv)}, xls sidecar: {len(new_sidecars)}, "
        f"PDF: {len(new_pdfs)}, pandoc(docx/rtf/odt/epub): {len(new_pandoc)}, pptx: {len(new_pptx)}, "
        f"md/txt: {len(new_md)}, code: {len(new_code)}, "
        f"MinerU(.doc/.ppt/html/图片): {len(new_others)}, 孤儿: {len(orphans)}"
    )

    if orphans:
        moves = stage_orphans(orphans, config.TARGET_DIR, config.ORPHANED_DIR, today)
        for src, dest in moves:
            logger.info(f"  Orphaned: {src} → {dest}")

    pruned = prune_empty_dirs(config.TARGET_DIR)
    if pruned:
        logger.info(f"Pruned {len(pruned)} empty dir(s):")
        for p in pruned:
            logger.info(f"  - {p.relative_to(config.TARGET_DIR)}")

    success = 0
    failed = 0

    def _run_local(label: str, files: list[Path], fn, what: str) -> None:
        """跑一组本地转换器：逐文件 try，单个失败不阻塞其余。"""
        nonlocal success, failed
        if not files:
            return
        logger.info(f"=== {label} {len(files)} 个 ===")
        for src in files:
            try:
                fn(src, logger)
                success += 1
            except Exception as exc:
                logger.error(f"  ✗ {src.name} {what}失败: {exc}")
                failed += 1

    # ── Stage 1: 本地 Office / Excel（无 MinerU 网络、不需 token）──
    _run_local("Excel 内容转换 (xlsx→openpyxl)", new_xlsx, process_xlsx, "xlsx 转换")
    _run_local("CSV/TSV 表格 (→md)", new_csv, process_csv, "csv 转换")
    _run_local("Excel sidecar (.xls)", new_sidecars, process_excel_sidecar, "sidecar ")
    _run_local("pandoc 转换 (docx/rtf/odt/epub)", new_pandoc, process_pandoc, "pandoc 转换")
    _run_local("PPT 转换 (pptx→python-pptx)", new_pptx, process_pptx, "pptx 转换")
    _run_local("Markdown/文本 直拷 (passthrough)", new_md, process_md_passthrough, "passthrough ")
    _run_local("代码/结构化文本 (code-fence)", new_code, process_code_passthrough, "code ")

    # ── Stage 2: PDF（先 pymupdf4llm，稀疏则 fallback 到 MinerU vlm）
    mineru_fallback_files: list[Path] = []
    if new_pdfs:
        logger.info(f"=== PDF 本地转换尝试 {len(new_pdfs)} 个 ===")
        for src in new_pdfs:
            try:
                md_content, pages, img_dir = convert_with_pymupdf4llm(src, logger)
            except Exception as exc:
                logger.warning(
                    f"  ⚠ {src.name}: pymupdf4llm 出错 ({exc})，转 MinerU fallback"
                )
                mineru_fallback_files.append(src)
                continue
            # 字符密度判断：太稀 → 多半扫描件 → fallback（丢弃已抽的图，交给 MinerU）
            density = (len(md_content) / pages) if pages else 0
            if density < config.PYMUPDF4LLM_MIN_CHARS_PER_PAGE:
                logger.warning(
                    f"  ⚠ {src.name}: pymupdf4llm 输出稀疏 "
                    f"({int(density)} chars/page < {config.PYMUPDF4LLM_MIN_CHARS_PER_PAGE})，"
                    f"判定扫描件，fallback MinerU vlm"
                )
                if img_dir:
                    shutil.rmtree(img_dir, ignore_errors=True)
                mineru_fallback_files.append(src)
                continue
            try:
                process_pymupdf4llm_result(src, md_content, pages, img_dir, logger)
                success += 1
            except Exception as exc:
                logger.error(f"  ✗ {src.name} pymupdf4llm 落地失败: {exc}")
                failed += 1

    # ── Stage 3: MinerU 收尾（fallback PDF + 其他非 PDF 文件）─────
    mineru_files = mineru_fallback_files + new_others

    # PDF fallback 用 vlm profile；其他文件按扩展名映射
    def _profile_for_mineru(p: Path) -> str:
        if p.suffix.lower() == ".pdf":
            return "vlm"
        return config.profile_name_for(p)

    def _print_summary(mineru_count: int) -> None:
        logger.info("=== Summary ===")
        logger.info(f"  Scanned: {len(candidates) + len(skipped)}")
        logger.info(f"  Skipped: {len(skipped)}")
        logger.info(f"  Excel 内容 (xlsx): {len(new_xlsx)}")
        logger.info(f"  CSV/TSV: {len(new_csv)}")
        logger.info(f"  Excel sidecar (.xls): {len(new_sidecars)}")
        logger.info(f"  pandoc (docx/rtf/odt/epub): {len(new_pandoc)}")
        logger.info(f"  PPT (pptx→python-pptx): {len(new_pptx)}")
        logger.info(f"  Markdown/文本 直拷: {len(new_md)}")
        logger.info(f"  代码/结构化文本: {len(new_code)}")
        logger.info(f"  pymupdf4llm 直出: {len(new_pdfs) - len(mineru_fallback_files)}")
        logger.info(f"  MinerU 处理: {mineru_count}")
        logger.info(f"    其中 PDF fallback: {len(mineru_fallback_files)}")
        logger.info(f"    其中 .doc/.ppt/html/图片: {len(new_others)}")
        logger.info(f"  Orphaned: {len(orphans)}")
        logger.info(f"  Success: {success}")
        logger.info(f"  Failed: {failed}")
        report_skipped(skipped, config.SOURCE_DIR, logger)
        logger.info("=== Sync finished ===")

    if not mineru_files:
        _print_summary(0)
        return 0 if failed == 0 else 1

    # Token 检查只在真正要调 MinerU 时做（sidecar / pymupdf4llm 不需要）
    if not config.MINERU_TOKEN:
        logger.error(
            f"需要 MinerU 转 {len(mineru_files)} 个文件但 MINERU_TOKEN 未设置；"
            f"sidecar 与 pymupdf4llm 已完成的部分保留"
        )
        return 1
    valid, msg, days = validate_token(config.MINERU_TOKEN)
    if not valid:
        logger.error(f"Token 检查失败: {msg}")
        return 1
    logger.info(f"Token: {msg}")
    if days < 14:
        logger.warning(f"⚠️  Token 仅剩 {days} 天，建议尽快更新")

    client = MinerUClient(config.MINERU_TOKEN, logger)
    files_by_profile: dict[str, list[Path]] = {}
    for src in mineru_files:
        files_by_profile.setdefault(_profile_for_mineru(src), []).append(src)

    batch_counter = 0
    for profile_name, profile_files in files_by_profile.items():
        profile = config.MINERU_PROFILES[profile_name]
        for start in range(0, len(profile_files), config.MAX_BATCH_SIZE):
            batch = profile_files[start:start + config.MAX_BATCH_SIZE]
            batch_counter += 1
            logger.info(
                f"=== Batch {batch_counter} | profile={profile_name} | "
                f"{len(batch)} files (streaming) ==="
            )
            try:
                for result in client.convert_batch_streaming(
                    batch,
                    language=config.MINERU_LANGUAGE,
                    **profile,
                ):
                    src = result.source_file
                    if result.success:
                        try:
                            process_result(src, result, logger)
                            success += 1
                        except Exception as exc:
                            logger.error(f"  ✗ {src.name} 落地失败: {exc}")
                            failed += 1
                    else:
                        logger.error(f"  ✗ {src.name}: {result.error}")
                        failed += 1
            except Exception as exc:
                logger.error(f"批次启动失败: {exc}")
                failed += len(batch)
                continue

    _print_summary(len(mineru_files))
    return 0 if failed == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
